import io
import os
import uuid
import logging
from dataclasses import dataclass, field

from app.config import settings

logger = logging.getLogger(__name__)


@dataclass
class ParseResult:
    text: str = ""
    images: list[dict] = field(default_factory=list)


async def parse_file(content: bytes, content_type: str) -> ParseResult:
    """Extract text and images from uploaded files."""
    try:
        if content_type == "application/pdf":
            return _parse_pdf(content)
        elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return _parse_pptx(content)
        elif content_type == "text/plain":
            return ParseResult(text=content.decode("utf-8", errors="replace"))
        else:
            return ParseResult()
    except Exception as e:
        logger.error(f"File parsing failed: {e}")
        return ParseResult()


def _parse_pdf(content: bytes) -> ParseResult:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    texts = []
    images = []

    os.makedirs(settings.upload_dir, exist_ok=True)

    for page_num, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            texts.append(text)

        # Extract embedded images
        try:
            for image in page.images:
                img_filename = f"{uuid.uuid4()}{os.path.splitext(image.name)[1] or '.png'}"
                img_path = os.path.join(settings.upload_dir, img_filename)
                with open(img_path, "wb") as f:
                    f.write(image.data)
                images.append({
                    "filename": img_filename,
                    "file_path": img_path,
                    "content_type": f"image/{os.path.splitext(image.name)[1].lstrip('.') or 'png'}",
                    "source_page": page_num + 1,
                })
        except Exception as e:
            logger.warning(f"Failed to extract images from PDF page {page_num + 1}: {e}")

    return ParseResult(text="\n\n".join(texts), images=images)


def _parse_pptx(content: bytes) -> ParseResult:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    texts = []
    images = []

    os.makedirs(settings.upload_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)

            # Extract embedded images
            if hasattr(shape, "image"):
                try:
                    img_blob = shape.image.blob
                    img_ct = shape.image.content_type or "image/png"
                    ext = img_ct.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    img_filename = f"{uuid.uuid4()}.{ext}"
                    img_path = os.path.join(settings.upload_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(img_blob)
                    images.append({
                        "filename": img_filename,
                        "file_path": img_path,
                        "content_type": img_ct,
                        "source_page": slide_num + 1,
                    })
                except Exception as e:
                    logger.warning(f"Failed to extract image from slide {slide_num + 1}: {e}")

    return ParseResult(text="\n\n".join(texts), images=images)
